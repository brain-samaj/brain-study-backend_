from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from app.ai.extractors.base import BaseExtractor
from app.ai.extractors.base import ExtractionResult


class PptxExtractor(BaseExtractor):
    SUPPORTED_EXTENSIONS = {".pptx"}

    def supports(self, suffix: str) -> bool:
        return suffix.lower() in self.SUPPORTED_EXTENSIONS

    def extract(
        self,
        source: Path,
    ) -> ExtractionResult:

        presentation = Presentation(source)

        slides: list[str] = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slides.append(text)

        return ExtractionResult(
            text="\n".join(slides),
            page_count=len(presentation.slides),
            metadata={
                "type": "pptx",
            },
        )
