from __future__ import annotations

from pathlib import Path

import fitz
import docx
from pptx import Presentation
from PIL import Image
import pytesseract


class DocumentProcessor:

    async def extract_text(
        self,
        path: str,
    ) -> str:

        suffix = Path(path).suffix.lower()

        if suffix == ".pdf":
            return self._pdf(path)

        if suffix == ".docx":
            return self._docx(path)

        if suffix == ".pptx":
            return self._pptx(path)

        if suffix == ".txt":
            return self._txt(path)

        if suffix in [
            ".jpg",
            ".jpeg",
            ".png",
            ".webp",
        ]:
            return self._image(path)

        raise ValueError("Unsupported file")


    def _pdf(
        self,
        path,
    ):

        document = fitz.open(path)

        text = []

        for page in document:
            text.append(page.get_text())

        return "\n".join(text)


    def _docx(
        self,
        path,
    ):

        document = docx.Document(path)

        return "\n".join(
            p.text
            for p in document.paragraphs
        )


    def _pptx(
        self,
        path,
    ):

        prs = Presentation(path)

        lines = []

        for slide in prs.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):
                    lines.append(shape.text)

        return "\n".join(lines)


    def _txt(
        self,
        path,
    ):

        with open(
            path,
            encoding="utf-8",
        ) as file:

            return file.read()


    def _image(
        self,
        path,
    ):

        image = Image.open(path)

        return pytesseract.image_to_string(image)

