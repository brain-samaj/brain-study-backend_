from __future__ import annotations

import json
import time
from pathlib import Path
from uuid import UUID

from app.ai.client import AIClient

from app.modules.knowledge_engine.models import KnowledgeSource
from app.modules.knowledge_engine.repository import KnowledgeRepository
from app.modules.knowledge_engine.schemas import (
    GlossaryItem,
    KnowledgeCreate,
    KnowledgeTopic,
    KnowledgeUpdate,
    LearningObjective,
    SampleQuestion,
)

from app.modules.study_materials.models import ProcessingStatus
from app.modules.study_materials.repository import (
    StudyMaterialRepository,
)


class KnowledgeEngineService:
    """
    Brain Study Knowledge Engine

    This service is the ONLY place that talks to an AI
    for learning-content generation.

    The generated KnowledgeSource becomes the permanent
    learning database for:

    • Study Guide
    • Flashcards
    • Smart Study
    • Exams
    • AI Tutor
    • Future learning features

    Unless the student presses "Regenerate",
    the same knowledge is reused forever.
    """

    def __init__(
        self,
        *,
        repository: KnowledgeRepository,
        study_material_repository: StudyMaterialRepository,
        ai_client: AIClient,
    ) -> None:

        self.repository = repository
        self.study_material_repository = study_material_repository
        self.ai = ai_client

    async def _load_material(
        self,
        material_id: UUID,
    ):

        material = await self.study_material_repository.get(
            material_id
        )

        if material is None:
            raise ValueError(
                "Study material not found."
            )

        return material

    async def _mark_processing(
        self,
        material,
    ):

        material.processing_status = (
            ProcessingStatus.PROCESSING
        )

        await self.study_material_repository._db.commit()
        await self.study_material_repository._db.refresh(
            material
        )

    async def _mark_ready(
        self,
        material,
    ):

        material.processing_status = (
            ProcessingStatus.READY
        )

        material.extraction_error = None

        await self.study_material_repository._db.commit()
        await self.study_material_repository._db.refresh(
            material
        )

    async def _mark_failed(
        self,
        material,
        error: Exception,
    ):

        material.processing_status = (
            ProcessingStatus.FAILED
        )

        material.extraction_error = str(error)

        await self.study_material_repository._db.commit()
        await self.study_material_repository._db.refresh(
            material
        )

    async def _extract_text(
        self,
        material,
    ) -> str:

        if (
            material.extracted_text
            and material.extracted_text.strip()
        ):
            return material.extracted_text

        path = Path(material.storage_path)

        if not path.exists():
            raise FileNotFoundError(
                f"Missing uploaded file: {path}"
            )

        suffix = path.suffix.lower()

        text = ""

        if suffix in {".txt", ".md"}:

            text = path.read_text(
                encoding="utf-8",
                errors="ignore",
            )

        elif suffix == ".pdf":

            from pypdf import PdfReader

            reader = PdfReader(str(path))

            pages = [
                page.extract_text() or ""
                for page in reader.pages
            ]

            text = "\n".join(pages)

            material.page_count = len(reader.pages)

        elif suffix == ".docx":

            from docx import Document

            document = Document(str(path))

            text = "\n".join(
                p.text
                for p in document.paragraphs
            )

        elif suffix == ".pptx":

            from pptx import Presentation

            presentation = Presentation(str(path))

            slides = []

            for slide in presentation.slides:

                for shape in slide.shapes:

                    if hasattr(shape, "text"):
                        slides.append(shape.text)

            text = "\n".join(slides)

            material.page_count = len(
                presentation.slides
            )

        elif suffix in {
            ".png",
            ".jpg",
            ".jpeg",
            ".bmp",
            ".gif",
            ".webp",
        }:

            text = ""

        else:

            raise RuntimeError(
                f"Unsupported file format: {suffix}"
            )

        text = text.strip()

        material.extracted_text = text

        material.word_count = (
            len(text.split())
            if text
            else 0
        )

        await self.study_material_repository._db.commit()

        await self.study_material_repository._db.refresh(
            material
        )

        return text

    # ==========================================================
    # BUILD KNOWLEDGE
    # ==========================================================

    async def build_from_material(
        self,
        *,
        material_id: UUID,
    ) -> KnowledgeSource:

        existing = await self.repository.get_by_material(
            material_id
        )

        if existing is not None:
            return existing

        material = await self._load_material(
            material_id
        )

        await self._mark_processing(
            material
        )

        try:

            extracted_text = await self._extract_text(
                material
            )

            if not extracted_text.strip():
                raise ValueError(
                    "No readable text extracted."
                )

            started = time.perf_counter()

            prompt = f"""
You are the Brain Study Knowledge Engine.

Your ONLY responsibility is to convert learning material into an intelligent structured knowledge database.

DO NOT teach.

DO NOT explain.

DO NOT summarize like a teacher.

Instead, analyse the material deeply.

Return ONLY valid JSON.

The JSON MUST follow this schema exactly.

{
  "title": "",
  "summary": "",

  "analysis": {
    "subject": "",

    "teaching_style": "",

    "explanation_style": "",

    "example_density": "",

    "reasoning_depth": "",

    "learning_styles": [],

    "best_teaching_methods": [],

    "recommended_learning_order": [],

    "common_student_mistakes": [],

    "real_world_applications": [],

    "prerequisites": [],

    "keywords": [],

    "important_terms": [],

    "needs_worked_examples": false,

    "needs_real_life_examples": false,

    "needs_visual_explanations": false,

    "needs_step_by_step": false,

    "needs_definitions_first": true,

    "needs_classification": false,

    "needs_comparison_tables": false,

    "needs_timelines": false,

    "needs_mnemonics": false,

    "requires_formulae": false,

    "requires_calculations": false,

    "requires_tables": false,

    "requires_diagrams": false,

    "requires_code": false,

    "requires_memorization": false
  },

  "topics": [
    {
      "title": "",
      "content": "",
      "keywords": [],
      "difficulty": "Easy | Medium | Hard"
    }
  ],

  "glossary": [
    {
      "term": "",
      "definition": ""
    }
  ],

  "learning_objectives": [
    {
      "objective": ""
    }
  ],

  "key_points": [],

  "sample_questions": [
    {
      "question": "",
      "answer": ""
    }
  ]
}

Rules

 Split large chapters into logical topics.

 Every topic must be self-contained.

 Include formulas exactly.

 Preserve mathematical symbols.

 Preserve chemical equations.

 Preserve physics equations.

 Preserve programming syntax.

 Preserve lists.

 Preserve tables as markdown if necessary.

 Include important definitions.

 Include important keywords.

 Never invent information.

Educational Analysis Rules

Determine:

- best teaching style
- reasoning depth
- example density
- learning order
- common student mistakes
- prerequisites
- important terms
- real-world applications

Determine whether the subject needs:

- formulas
- calculations
- diagrams
- timelines
- code examples
- classification tables
- worked examples
- mnemonics
- step-by-step teaching

These decisions will later be used by TeacherAI.

TeacherAI must NEVER re-analyse the material.

Study Material:

{extracted_text}
"""

            response = await self.ai.generate_json(
                prompt=prompt,
                temperature=0.15,
            )

            analysis = response.get("analysis", {})

            analysis.setdefault("subject", response.get("title", material.title))
            analysis.setdefault("teaching_style", "Adaptive")
            analysis.setdefault("explanation_style", "Simple")
            analysis.setdefault("example_density", "Moderate")
            analysis.setdefault("reasoning_depth", "Intermediate")

            analysis.setdefault("needs_worked_examples", False)
            analysis.setdefault("needs_real_life_examples", False)
            analysis.setdefault("needs_visual_explanations", False)
            analysis.setdefault("needs_step_by_step", False)
            analysis.setdefault("needs_definitions_first", True)
            analysis.setdefault("needs_classification", False)
            analysis.setdefault("needs_comparison_tables", False)
            analysis.setdefault("needs_timelines", False)
            analysis.setdefault("needs_mnemonics", False)

            analysis.setdefault("requires_formulae", False)
            analysis.setdefault("requires_calculations", False)
            analysis.setdefault("requires_tables", False)
            analysis.setdefault("requires_diagrams", False)
            analysis.setdefault("requires_code", False)
            analysis.setdefault("requires_memorization", False)

            analysis.setdefault("keywords", [])
            analysis.setdefault("important_terms", [])
            analysis.setdefault("prerequisites", [])
            analysis.setdefault("learning_styles", ["mixed"])
            analysis.setdefault("best_teaching_methods", [])
            analysis.setdefault("common_student_mistakes", [])
            analysis.setdefault("real_world_applications", [])
            analysis.setdefault("recommended_learning_order", [])

            response["analysis"] = analysis

            elapsed = int(
                (
                    time.perf_counter()
                    - started
                )
                * 1000
            )

            knowledge = await self.repository.create(

                KnowledgeCreate(

                    material_id=material.id,

                    title=response.get(
                        "title",
                        material.title,
                    ),

                    summary=response.get(
                        "summary",
                        "",
                    ),

                    knowledge=response,

                    topics=[
                        KnowledgeTopic(**item)
                        for item in response.get(
                            "topics",
                            [],
                        )
                    ],

                    glossary=[
                        GlossaryItem(**item)
                        for item in response.get(
                            "glossary",
                            [],
                        )
                    ],

                    learning_objectives=[
                        LearningObjective(**item)
                        for item in response.get(
                            "learning_objectives",
                            [],
                        )
                    ],

                    key_points=response.get(
                        "key_points",
                        [],
                    ),

                    sample_questions=[
                        SampleQuestion(**item)
                        for item in response.get(
                            "sample_questions",
                            [],
                        )
                    ],

                    total_tokens=0,

                    ai_provider="Brain AI",

                    ai_model="Knowledge Engine",

                    processing_time_ms=elapsed,

                    is_cached=False,
                )
            )

            await self._mark_ready(
                material
            )

            return knowledge

        except Exception as exc:

            await self._mark_failed(
                material,
                exc,
            )

            raise


    # ==========================================================
    # REFRESH KNOWLEDGE
    # ==========================================================

    async def refresh(
        self,
        *,
        material_id: UUID,
    ) -> KnowledgeSource:
        """
        Deletes the existing knowledge and rebuilds it.
        """

        await self.repository.delete(
            material_id
        )

        return await self.build_from_material(
            material_id=material_id,
        )

    # ==========================================================
    # GET KNOWLEDGE
    # ==========================================================

    async def get(
        self,
        *,
        material_id: UUID,
    ) -> KnowledgeSource | None:

        return await self.repository.get_by_material(
            material_id
        )

    async def get_or_build(
        self,
        *,
        material_id: UUID,
    ) -> KnowledgeSource:

        knowledge = await self.repository.get_by_material(
            material_id
        )

        if knowledge is not None:
            return knowledge

        return await self.build_from_material(
            material_id=material_id,
        )

    # ==========================================================
    # UPDATE KNOWLEDGE
    # ==========================================================

    async def update(
        self,
        *,
        knowledge: KnowledgeSource,
        payload: KnowledgeUpdate,
    ) -> KnowledgeSource:

        return await self.repository.update(
            knowledge,
            payload,
        )

    # ==========================================================
    # DELETE KNOWLEDGE
    # ==========================================================

    async def delete(
        self,
        *,
        material_id: UUID,
    ) -> bool:

        return await self.repository.delete(
            material_id
        )

    # ==========================================================
    # EXISTS
    # ==========================================================

    async def exists(
        self,
        *,
        material_id: UUID,
    ) -> bool:

        return await self.repository.exists(
            material_id
        )

    # ==========================================================
    # EXPORT JSON
    # ==========================================================

    async def export_json(
        self,
        *,
        material_id: UUID,
    ) -> dict:

        knowledge = await self.get_or_build(
            material_id=material_id,
        )

        return {
            "id": str(
                knowledge.id
            ),
            "material_id": str(
                knowledge.material_id
            ),
            "status": knowledge.status.value,
            "title": knowledge.title,
            "summary": knowledge.summary,
            "knowledge": knowledge.knowledge,
            "topics": knowledge.topics,
            "glossary": knowledge.glossary,
            "learning_objectives":
                knowledge.learning_objectives,
            "key_points":
                knowledge.key_points,
            "sample_questions":
                knowledge.sample_questions,
            "total_tokens":
                knowledge.total_tokens,
            "provider":
                knowledge.ai_provider,
            "model":
                knowledge.ai_model,
            "processing_time_ms":
                knowledge.processing_time_ms,
            "cached":
                knowledge.is_cached,
        }

    # ==========================================================
    # EXPORT SUMMARY
    # ==========================================================

    async def export_summary(
        self,
        *,
        material_id: UUID,
    ) -> dict:

        knowledge = await self.get_or_build(
            material_id=material_id,
        )

        return {

            "title": knowledge.title,

            "summary": knowledge.summary,

            "topics": knowledge.topics,

            "key_points":
                knowledge.key_points,

            "learning_objectives":
                knowledge.learning_objectives,

            "glossary":
                knowledge.glossary,

            "sample_questions":
                knowledge.sample_questions,
        }


    # ==========================================================
    # PROCESSING STATUS
    # ==========================================================

    async def get_processing_status(
        self,
        *,
        material_id: UUID,
    ) -> dict:

        material = await self.study_material_repository.get(
            material_id
        )

        if material is None:
            raise ValueError(
                "Study material not found."
            )

        knowledge = await self.repository.get_by_material(
            material_id
        )

        return {

            "material_id": str(material.id),

            "processing_status": (
                material.processing_status.value
                if material.processing_status
                else None
            ),

            "knowledge_status": (
                knowledge.status.value
                if knowledge
                else None
            ),

            "has_knowledge": knowledge is not None,

            "title": (
                knowledge.title
                if knowledge
                else None
            ),

            "word_count": material.word_count,

            "page_count": material.page_count,

            "error": material.extraction_error,
        }

    # ==========================================================
    # REBUILD KNOWLEDGE
    # ==========================================================

    async def rebuild(
        self,
        *,
        material_id: UUID,
    ) -> KnowledgeSource:

        material = await self._load_material(
            material_id
        )

        await self.repository.delete(
            material_id
        )

        material.processing_status = (
            ProcessingStatus.UPLOADING
        )

        material.extracted_text = ""

        material.word_count = 0

        material.page_count = None

        material.extraction_error = None

        await self.study_material_repository._db.commit()

        return await self.build_from_material(
            material_id=material_id,
        )

    # ==========================================================
    # SEARCH KNOWLEDGE
    # ==========================================================

    async def search(
        self,
        *,
        material_id: UUID,
        keyword: str,
    ) -> dict:

        knowledge = await self.get_or_build(
            material_id=material_id,
        )

        keyword = keyword.lower().strip()

        results = []

        for topic in knowledge.topics or []:

            if keyword in str(topic).lower():

                results.append(
                    {
                        "section": "topic",
                        "title": topic.get("title"),
                        "content": topic.get("content"),
                    }
                )

        for item in knowledge.glossary or []:

            if keyword in str(item).lower():

                results.append(
                    {
                        "section": "glossary",
                        "term": item.get("term"),
                        "definition": item.get("definition"),
                    }
                )

        for point in knowledge.key_points or []:

            if keyword in point.lower():

                results.append(
                    {
                        "section": "key_point",
                        "content": point,
                    }
                )

        if keyword in knowledge.summary.lower():

            results.append(
                {
                    "section": "summary",
                    "content": knowledge.summary,
                }
            )

        return {

            "keyword": keyword,

            "count": len(results),

            "results": results,
        }

    # ==========================================================
    # VALIDATE KNOWLEDGE
    # ==========================================================

    async def validate(
        self,
        *,
        material_id: UUID,
    ) -> bool:

        knowledge = await self.repository.get_by_material(
            material_id
        )

        if knowledge is None:
            return False

        if not knowledge.summary.strip():
            return False

        if not knowledge.topics:
            return False

        if not knowledge.knowledge:
            return False

        return True


    # ==========================================================
    # CACHE UTILITIES
    # ==========================================================

    async def clear_cache(
        self,
        *,
        material_id: UUID,
    ) -> None:
        """
        Removes the cached knowledge.

        The next request will regenerate it.
        """

        await self.repository.delete(
            material_id
        )

    async def regenerate(
        self,
        *,
        material_id: UUID,
        force: bool = False,
    ) -> KnowledgeSource:
        """
        Returns cached knowledge unless force=True.
        """

        if force:
            return await self.rebuild(
                material_id=material_id,
            )

        return await self.get_or_build(
            material_id=material_id,
        )

    # ==========================================================
    # NORMALIZATION HELPERS
    # ==========================================================

    def normalize_topic(
        self,
        topic: dict,
    ) -> dict:

        return {

            "title": (
                topic.get("title") or ""
            ).strip(),

            "content": (
                topic.get("content") or ""
            ).strip(),

            "keywords": topic.get(
                "keywords",
                [],
            ),

            "difficulty": (
                topic.get("difficulty")
                or "Medium"
            ),
        }

    def normalize_glossary(
        self,
        glossary: list,
    ) -> list:

        return [
            {
                "term": item.get(
                    "term",
                    "",
                ).strip(),

                "definition": item.get(
                    "definition",
                    "",
                ).strip(),
            }
            for item in glossary
        ]

    # ==========================================================
    # FUTURE MODULE HELPERS
    # ==========================================================

    async def get_topics(
        self,
        *,
        material_id: UUID,
    ) -> list:

        knowledge = await self.get_or_build(
            material_id=material_id,
        )

        return [
            self.normalize_topic(topic)
            for topic in (
                knowledge.topics or []
            )
        ]

    async def get_glossary(
        self,
        *,
        material_id: UUID,
    ) -> list:

        knowledge = await self.get_or_build(
            material_id=material_id,
        )

        return self.normalize_glossary(
            knowledge.glossary or []
        )

    async def get_learning_objectives(
        self,
        *,
        material_id: UUID,
    ) -> list:

        knowledge = await self.get_or_build(
            material_id=material_id,
        )

        return knowledge.learning_objectives or []

    async def get_key_points(
        self,
        *,
        material_id: UUID,
    ) -> list:

        knowledge = await self.get_or_build(
            material_id=material_id,
        )

        return knowledge.key_points or []

    async def get_sample_questions(
        self,
        *,
        material_id: UUID,
    ) -> list:

        knowledge = await self.get_or_build(
            material_id=material_id,
        )

        return knowledge.sample_questions or []

