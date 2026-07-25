from __future__ import annotations

from pathlib import Path

from docx import Document
from pptx import Presentation
from pypdf import PdfReader


class StudyMaterialExtractor:
    """
    Brain Study Universal Text Extractor.

    Supported formats
    -----------------
    • PDF
    • DOCX
    • PPTX
    • TXT
    • MD
    • PNG
    • JPG
    • JPEG
    • WEBP
    • BMP
    • GIF

    This class ONLY extracts text.

    It NEVER:
        - writes to the database
        - updates models
        - commits transactions
    """

    IMAGE_TYPES = {
        ".png",
        ".jpg",
        ".jpeg",
        ".bmp",
        ".gif",
        ".webp",
    }

    TEXT_TYPES = {
        ".txt",
        ".md",
    }

    DOCUMENT_TYPES = {
        ".pdf",
        ".docx",
        ".pptx",
    }

    SUPPORTED_TYPES = (
        IMAGE_TYPES
        | TEXT_TYPES
        | DOCUMENT_TYPES
    )

    # ==========================================================
    # PUBLIC
    # ==========================================================

    async def extract(
        self,
        file_path: str | Path,
    ) -> str:

        path = Path(file_path)

        if not path.exists():
            raise FileNotFoundError(path)

        suffix = path.suffix.lower()

        if suffix in self.TEXT_TYPES:
            return self._extract_plain_text(path)

        if suffix == ".pdf":
            return self._extract_pdf(path)

        if suffix == ".docx":
            return self._extract_docx(path)

        if suffix == ".pptx":
            return self._extract_pptx(path)

        if suffix in self.IMAGE_TYPES:
            return self._extract_image(path)

        raise RuntimeError(
            f"Unsupported file type: {suffix}"
        )

    # ==========================================================
    # TXT / MD
    # ==========================================================

    def _extract_plain_text(
        self,
        path: Path,
    ) -> str:

        return path.read_text(
            encoding="utf-8",
            errors="ignore",
        ).strip()

    # ==========================================================
    # PDF
    # ==========================================================

    def _extract_pdf(
        self,
        path: Path,
    ) -> str:

        reader = PdfReader(str(path))

        pages: list[str] = []

        for page in reader.pages:
            pages.append(page.extract_text() or "")

        return "\n".join(pages).strip()

    # ==========================================================
    # DOCX
    # ==========================================================

    def _extract_docx(
        self,
        path: Path,
    ) -> str:

        document = Document(str(path))

        paragraphs: list[str] = []

        for paragraph in document.paragraphs:

            text = paragraph.text.strip()

            if text:
                paragraphs.append(text)

        return "\n".join(paragraphs).strip()

    # ==========================================================
    # PPTX
    # ==========================================================

    def _extract_pptx(
        self,
        path: Path,
    ) -> str:

        presentation = Presentation(str(path))

        contents: list[str] = []

        for slide in presentation.slides:

            for shape in slide.shapes:

                if not hasattr(shape, "text"):
                    continue

                text = shape.text.strip()

                if text:
                    contents.append(text)

        return "\n".join(contents).strip()

    # ==========================================================
    # IMAGE OCR
    # ==========================================================

    def _extract_image(
        self,
        path: Path,
    ) -> str:

        try:
            from PIL import Image
            import pytesseract
        except ImportError as exc:
            raise RuntimeError(
                "Image OCR requires pillow and pytesseract."
            ) from exc

        image = Image.open(path)

        text = pytesseract.image_to_string(
            image,
            lang="eng",
        )

        return text.strip()

    # ==========================================================
    # VALIDATION
    # ==========================================================

    @classmethod
    def supports(
        cls,
        filename: str,
    ) -> bool:

        suffix = Path(filename).suffix.lower()

        return suffix in cls.SUPPORTED_TYPES

    @classmethod
    def supported_extensions(
        cls,
    ) -> list[str]:

        return sorted(cls.SUPPORTED_TYPES)
